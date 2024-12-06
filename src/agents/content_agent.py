from src.api.chromadb_api import add_document_to_chromadb
from PyPDF2 import PdfReader
from youtube_transcript_api import YouTubeTranscriptApi
from pptx import Presentation


class ContentIngestionAgent:

    def process_pdf(self, file):
        """
        Process a PDF file, extract text, split into chunks, and store it in the vector database.
        """
        try:
            reader = PdfReader(file)
            text_chunks = []

            # Split text into chunks
            for page_num, page in enumerate(reader.pages):
                text = page.extract_text()
                chunks = text.split("\n\n")  # Split by paragraph
                for chunk in chunks:
                    metadata = {
                        "file_name": file.name,
                        "type": "pdf",
                        "page_num": page_num + 1,
                        "unique_id": f"{file.name}_page_{page_num + 1}"
                    }
                    text_chunks.append({"content": chunk, "metadata": metadata})

            # Add pre-chunked content to ChromaDB
            return add_document_to_chromadb(chunks=text_chunks)
        except Exception as e:
            return f"Error processing PDF file '{file.name}': {e}"
        

    def process_youtube_video(self, video_url):
        """
        Extract transcript from a YouTube video and store it in the vector database.
        """
        try:
            video_id = video_url.split("v=")[-1]
            transcript = YouTubeTranscriptApi.get_transcript(video_id)
            content = " ".join([t['text'] for t in transcript])

            metadata = {
                "file_name": f"YouTube_{video_id}",
                "type": "youtube",
                "unique_id": video_id
            }

            # Add content and metadata to ChromaDB
            return add_document_to_chromadb(content=content, metadata=metadata)
        except Exception as e:
            return f"Error processing YouTube video: {e}"
      
    def process_pptx(self, file):
        """
        Extract text and images from a PowerPoint presentation and store it in the vector database.
        """
        try:
            presentation = Presentation(file)
            text_chunks = []

            for slide_idx, slide in enumerate(presentation.slides, start=1):
                for shape_idx, shape in enumerate(slide.shapes, start=1):
                    if shape.has_text_frame:
                        text = shape.text.strip()
                        metadata = {
                            "file_name": file.name,
                            "type": "pptx",
                            "slide_num": slide_idx,
                            "unique_id": f"{file.name}_slide_{slide_idx}_shape_{shape_idx}"
                        }
                        text_chunks.append({"content": text, "metadata": metadata})

            # Add pre-chunked content to ChromaDB
            return add_document_to_chromadb(chunks=text_chunks)
        except Exception as e:
            return f"Error processing PowerPoint file '{file.name}': {e}"